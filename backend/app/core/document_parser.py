"""
Multi-format document parser for RAG system
"""

import os
import hashlib
from pathlib import Path
from typing import List, Dict, Any, Optional
from datetime import datetime

# Document parsing libraries
import PyPDF2
import pdfplumber
from docx import Document
from pptx import Presentation
import markdown

from backend.app.core.config import settings
from backend.app.utils.logger import logger
from backend.app.models.schemas import DocumentMetadata

class DocumentParser:
    """Multi-format document parser"""
    
    def __init__(self):
        self.supported_extensions = settings.allowed_extensions
        
    def parse_document(self, file_path: str) -> Dict[str, Any]:
        """Parse document and extract text content"""
        try:
            file_path = Path(file_path)
            extension = file_path.suffix.lower()
            
            if extension not in self.supported_extensions:
                raise ValueError(f"Unsupported file format: {extension}")
            
            # Extract metadata
            metadata = self._extract_metadata(file_path)
            
            # Parse content based on file type
            if extension == '.pdf':
                content = self._parse_pdf(file_path)
            elif extension == '.docx':
                content = self._parse_docx(file_path)
            elif extension == '.pptx':
                content = self._parse_pptx(file_path)
            elif extension == '.txt':
                content = self._parse_txt(file_path)
            elif extension == '.md':
                content = self._parse_markdown(file_path)
            else:
                raise ValueError(f"Parser not implemented for {extension}")
            
            # Clean and preprocess content
            content = self._clean_text(content)
            
            return {
                'content': content,
                'metadata': metadata,
                'success': True
            }
            
        except Exception as e:
            logger.error(f"Error parsing document {file_path}: {str(e)}")
            return {
                'content': '',
                'metadata': None,
                'success': False,
                'error': str(e)
            }
    
    def _extract_metadata(self, file_path: Path) -> DocumentMetadata:
        """Extract document metadata"""
        stat = file_path.stat()
        
        return DocumentMetadata(
            filename=file_path.name,
            file_type=file_path.suffix.lower(),
            file_size=stat.st_size,
            upload_time=datetime.fromtimestamp(stat.st_mtime)
        )
    
    def _parse_pdf(self, file_path: Path) -> str:
        """Parse PDF document"""
        content = ""
        
        try:
            # Try with pdfplumber first (better text extraction)
            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages:
                    page_text = page.extract_text()
                    if page_text:
                        content += page_text + "\n"
        except Exception as e:
            logger.warning(f"pdfplumber failed for {file_path}, trying PyPDF2: {e}")
            
            # Fallback to PyPDF2
            try:
                with open(file_path, 'rb') as file:
                    pdf_reader = PyPDF2.PdfReader(file)
                    for page in pdf_reader.pages:
                        content += page.extract_text() + "\n"
            except Exception as e2:
                logger.error(f"Both PDF parsers failed for {file_path}: {e2}")
                raise e2
        
        return content
    
    def _parse_docx(self, file_path: Path) -> str:
        """Parse Word document"""
        doc = Document(file_path)
        content = ""
        
        # Extract paragraphs
        for paragraph in doc.paragraphs:
            content += paragraph.text + "\n"
        
        # Extract tables
        for table in doc.tables:
            for row in table.rows:
                row_text = []
                for cell in row.cells:
                    row_text.append(cell.text.strip())
                content += " | ".join(row_text) + "\n"
        
        return content
    
    def _parse_pptx(self, file_path: Path) -> str:
        """Parse PowerPoint presentation"""
        prs = Presentation(file_path)
        content = ""
        
        for slide_num, slide in enumerate(prs.slides, 1):
            content += f"\n--- Slide {slide_num} ---\n"
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    content += shape.text + "\n"
                
                # Extract table content if present
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_text = []
                        for cell in row.cells:
                            row_text.append(cell.text.strip())
                        content += " | ".join(row_text) + "\n"
        
        return content
    
    def _parse_txt(self, file_path: Path) -> str:
        """Parse text file"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                return file.read()
        except UnicodeDecodeError:
            # Try with different encoding
            with open(file_path, 'r', encoding='gbk') as file:
                return file.read()
    
    def _parse_markdown(self, file_path: Path) -> str:
        """Parse Markdown file"""
        with open(file_path, 'r', encoding='utf-8') as file:
            md_content = file.read()
        
        # Convert markdown to plain text
        html = markdown.markdown(md_content)
        # Simple HTML tag removal (for basic conversion)
        import re
        text = re.sub('<[^<]+?>', '', html)
        return text
    
    def _clean_text(self, text: str) -> str:
        """Clean and preprocess text"""
        if not text:
            return ""
        
        # Remove excessive whitespace
        text = ' '.join(text.split())
        
        # Remove special characters but keep basic punctuation
        import re
        text = re.sub(r'[^\w\s\.\,\!\?\;\:\-\(\)\[\]\"\']+', ' ', text)
        
        # Remove excessive newlines
        text = re.sub(r'\n+', '\n', text)
        
        return text.strip()
    
    def validate_file(self, file_path: str) -> Dict[str, Any]:
        """Validate file before processing"""
        try:
            file_path = Path(file_path)
            
            # Check if file exists
            if not file_path.exists():
                return {'valid': False, 'error': 'File does not exist'}
            
            # Check file extension
            if file_path.suffix.lower() not in self.supported_extensions:
                return {'valid': False, 'error': f'Unsupported file format: {file_path.suffix}'}
            
            # Check file size
            file_size_mb = file_path.stat().st_size / (1024 * 1024)
            if file_size_mb > settings.max_file_size_mb:
                return {'valid': False, 'error': f'File too large: {file_size_mb:.1f}MB > {settings.max_file_size_mb}MB'}
            
            return {'valid': True}
            
        except Exception as e:
            return {'valid': False, 'error': str(e)}

# Global parser instance
document_parser = DocumentParser()
